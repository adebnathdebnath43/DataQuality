from typing import List, Dict, Any, Optional
import json
import io
import datetime
from pathlib import Path
from app.services.s3 import S3Service
from app.services.bedrock import BedrockService

# Import text extraction libraries
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

LOG_PATH = Path(__file__).resolve().parents[2] / "debug_absolute.log"


class MetadataService:
    def __init__(self):
        self.s3_service = S3Service()
        self.bedrock_service = BedrockService()

    def _extract_text(self, content: bytes, file_ext: str) -> str:
        """Extract text from various file formats"""
        file_ext = file_ext.lower()
        
        if file_ext == 'pdf':
            if not pypdf:
                return "Error: pypdf library not installed"
            try:
                pdf_file = io.BytesIO(content)
                reader = pypdf.PdfReader(pdf_file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            except Exception as e:
                return f"Error extracting PDF text: {str(e)}"

        elif file_ext in ['docx', 'doc']:
            if not docx:
                return "Error: python-docx library not installed"
            try:
                doc_file = io.BytesIO(content)
                doc = docx.Document(doc_file)
                return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            except Exception as e:
                return f"Error extracting DOCX text: {str(e)}"

        elif file_ext in ['pptx', 'ppt']:
            if not Presentation:
                return "Error: python-pptx library not installed"
            try:
                ppt_file = io.BytesIO(content)
                prs = Presentation(ppt_file)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            except Exception as e:
                return f"Error extracting PPTX text: {str(e)}"

        else:
            # Assume text for other formats
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                return f"Error: Binary file {file_ext} not supported for text extraction"

    def _cosine_similarity(self, vec1: List[float], vec2: List[float]) -> float:
        """Calculate cosine similarity between two vectors"""
        try:
            import math
            
            # Ensure vectors are same length
            if len(vec1) != len(vec2):
                return 0.0
            
            # Calculate dot product
            dot_product = sum(a * b for a, b in zip(vec1, vec2))
            
            # Calculate magnitudes
            magnitude1 = math.sqrt(sum(a * a for a in vec1))
            magnitude2 = math.sqrt(sum(b * b for b in vec2))
            
            # Avoid division by zero
            if magnitude1 == 0 or magnitude2 == 0:
                return 0.0
            
            # Calculate cosine similarity
            return dot_product / (magnitude1 * magnitude2)
        except Exception as e:
            self._log(f"Error calculating cosine similarity: {str(e)}")
            return 0.0
    
    def _log(self, msg: str):
        try:
            with open(LOG_PATH, "a", encoding="utf-8") as f:
                f.write(f"{datetime.datetime.now()}: {msg}\n")
        except Exception as e:
            print(f"Logging failed: {e}")
        print(msg)

    def _parse_flexible_date(self, value: Any) -> Optional[datetime.datetime]:
        """Parse common ISO-ish or US-style date strings into datetime."""
        try:
            if isinstance(value, datetime.datetime):
                return value
            if value is None:
                return None
            text = str(value).strip()
            if not text:
                return None
            normalized = text.replace("Z", "+00:00")
            # Try ISO first
            try:
                return datetime.datetime.fromisoformat(normalized)
            except Exception:
                pass
            # Try common formats
            for fmt in [
                "%Y-%m-%d",
                "%Y/%m/%d",
                "%m/%d/%Y",
                "%m/%d/%y",
                "%Y-%m-%d %H:%M:%S",
                "%Y/%m/%d %H:%M:%S",
                "%m/%d/%Y %H:%M:%S",
                "%m/%d/%y %H:%M:%S",
            ]:
                try:
                    return datetime.datetime.strptime(text, fmt)
                except Exception:
                    continue
        except Exception:
            return None
        return None

    def _validate_dimensions(self, dimensions: Dict[str, Any]) -> Dict[str, Any]:
        """Validate and ensure all 17 dimensions are present with valid scores"""
        required_dimensions = [
            "Accuracy", "Completeness", "Consistency", "Timeliness", "Validity",
            "Uniqueness", "Reliability", "Relevance", "Accessibility", "Precision",
            "Integrity", "Conformity", "Interpretability", "Traceability",
            "Credibility", "Fitness_for_Use", "Value"
        ]

        # Map incoming keys case-insensitively to the required names
        normalized_map = {dim.lower(): dim for dim in required_dimensions}
        incoming = {}
        for k, v in (dimensions or {}).items():
            target = normalized_map.get(str(k).lower())
            if target:
                incoming[target] = v
        
        validated = {}
        for dim in required_dimensions:
            if dim in incoming and isinstance(incoming[dim], dict):
                score = incoming[dim].get("score", 50)
                evidence = incoming[dim].get("evidence", "Not assessed")
                # Ensure score is valid (0-100)
                score = max(0, min(100, int(score)))
                validated[dim] = {"score": score, "evidence": evidence}
            else:
                # Default for missing dimensions
                validated[dim] = {"score": 50, "evidence": "Dimension not assessed by LLM"}
        
        return validated

    def _cosine_similarity(self, v1: List[float], v2: List[float]) -> float:
        """Calculate cosine similarity between two vectors"""
        if not v1 or not v2 or len(v1) != len(v2):
            return 0.0
        
        dot_product = sum(a*b for a,b in zip(v1, v2))
        magnitude1 = sum(a*a for a in v1) ** 0.5
        magnitude2 = sum(b*b for b in v2) ** 0.5
        
        if magnitude1 == 0 or magnitude2 == 0:
            return 0.0
            
        return dot_product / (magnitude1 * magnitude2)

    def find_duplicates(self, target_embedding: List[float], all_results: List[Dict[str, Any]], threshold: float = 0.95) -> List[Dict[str, Any]]:
        """Find duplicate documents based on embedding similarity"""
        duplicates = []
        for result in all_results:
            if 'embedding' in result and result['embedding']:
                similarity = self._cosine_similarity(target_embedding, result['embedding'])
                if similarity >= threshold:
                    duplicates.append({
                        "file_name": result.get('file_name', 'Unknown'),
                        "similarity": round(similarity * 100, 2),
                        "file_key": result.get('file_key', '')
                    })
        
        # Sort by similarity (highest first)
        duplicates.sort(key=lambda x: x['similarity'], reverse=True)
        return duplicates

    def _metadata_similarity(self, doc_type1: str, doc_type2: str, topics1: Any, topics2: Any, key_terms1: Any, key_terms2: Any) -> float:
        """Combined similarity on doc type (exact match) and Jaccard over topics/key terms."""
        try:
            import re

            def tokenize_list(value: Any) -> set:
                if isinstance(value, list):
                    items = value
                elif value is None:
                    items = []
                else:
                    items = [value]
                text = " ".join([str(x) for x in items]).lower()
                return set([t for t in re.split(r"\W+", text) if t])

            doc_type_score = 1.0 if (doc_type1 or "").strip().lower() == (doc_type2 or "").strip().lower() and doc_type1 else 0.0
            s1 = tokenize_list(topics1) | tokenize_list(key_terms1)
            s2 = tokenize_list(topics2) | tokenize_list(key_terms2)

            if not s1 or not s2:
                jaccard = 0.0
            else:
                inter = len(s1 & s2)
                union = len(s1 | s2)
                jaccard = inter / union if union else 0.0

            # Weighted blend: doc type 0.5, topics/key terms 0.5
            return 0.5 * doc_type_score + 0.5 * jaccard
        except Exception as e:
            self._log(f"Metadata similarity failed: {str(e)}")
            return 0.0

    async def process_files(self, bucket: str, file_keys: List[str], region: str = None, access_key: str = None, secret_key: str = None, role_arn: str = None, model_id: str = None) -> List[Dict[str, Any]]:
        self._log(f"Processing {len(file_keys)} files from bucket {bucket}")
        results = []
        file_analyses = []  # Store full analysis for each file
        
        for key in file_keys:
            self._log(f"Starting processing for file: {key}")
            try:
                # Skip folders
                if key.endswith('/'):
                    self._log(f"Skipping folder: {key}")
                    continue

                # Determine file type
                file_ext = key.split('.')[-1] if '.' in key else ''
                self._log(f"File extension: {file_ext}")

                # Fetch S3 object metadata for upload date
                obj_meta = await self.s3_service.get_object_metadata(bucket, key, region, access_key, secret_key, role_arn)
                upload_dt = obj_meta.get("last_modified")
                upload_date_iso = upload_dt.isoformat() if upload_dt else None
                upload_age_days = None
                if upload_dt:
                    # Normalize to naive UTC for difference calculation
                    try:
                        upload_age_days = (datetime.datetime.utcnow() - upload_dt.replace(tzinfo=None)).days
                    except Exception as age_err:
                        self._log(f"Failed to compute upload_age_days for {key}: {str(age_err)}")
                
                # 1. Read file content (binary mode)
                self._log(f"Reading file from S3: {key}")
                content_bytes = await self.s3_service.read_file(bucket, key, region, access_key, secret_key, role_arn, binary=True)
                self._log(f"Read {len(content_bytes)} bytes")
                
                # 2. Extract text
                self._log("Extracting text...")
                text_content = self._extract_text(content_bytes, file_ext)
                self._log(f"Extracted text length: {len(text_content)}")
                
                if not text_content or text_content.startswith("Error:"):
                    error_msg = text_content or "Empty content"
                    self._log(f"Text extraction failed: {error_msg}")
                    results.append({
                        "file_key": key,
                        "status": "error",
                        "error": error_msg
                    })
                    file_analyses.append({
                        "file_key": key,
                        "file_name": key.split('/')[-1],
                        "status": "error",
                        "error": error_msg,
                        "processed_at": datetime.datetime.utcnow().isoformat() + "Z"
                    })
                    continue

                # 3. Analyze with Bedrock
                file_name = key.split('/')[-1]
                self._log(f"Analyzing with Bedrock model: {model_id}")
                analysis = self.bedrock_service.analyze_content(
                    content=text_content, 
                    file_name=file_name,
                    model_id=model_id,
                    region=region,
                    access_key=access_key,
                    secret_key=secret_key,
                    role_arn=role_arn
                )
                self._log("Bedrock analysis complete")
                
                # 3.5 Generate Embedding
                embedding = []
                try:
                    # Create text to embed: ALL metadata fields + summary + context
                    summary = analysis.get("summary", "")
                    context = analysis.get("context", "")
                    document_type = analysis.get("document_type", "")
                    
                    # Extract all metadata fields
                    metadata = analysis.get("metadata", {})
                    metadata_parts = []
                    
                    for key, value in metadata.items():
                        if isinstance(value, str):
                            metadata_parts.append(f"{key}: {value}")
                        elif isinstance(value, list):
                            metadata_parts.append(f"{key}: {', '.join(str(v) for v in value)}")
                        else:
                            metadata_parts.append(f"{key}: {str(value)}")
                    
                    metadata_text = "\n".join(metadata_parts)
                    
                    # Combine all content for comprehensive embedding
                    text_to_embed = f"""Document Type: {document_type}

Summary: {summary}

Context: {context}

Metadata:
{metadata_text}

Full Content (truncated): {text_content[:2000]}"""
                    
                    self._log(f"Generating embedding for comprehensive content (text length: {len(text_to_embed)})...")
                    embedding = self.bedrock_service.get_embedding(
                        text=text_to_embed,
                        region=region,
                        access_key=access_key,
                        secret_key=secret_key,
                        role_arn=role_arn
                    )
                    self._log(f"Embedding generated. Length: {len(embedding)}, First 3 values: {embedding[:3] if embedding else 'EMPTY'}")
                except Exception as embed_err:
                    self._log(f"Embedding generation failed: {str(embed_err)}")
                    import traceback
                    self._log(traceback.format_exc())

                # 3.6 Validate and extract dimensions
                dimensions = {}
                recommended_action = (analysis.get("recommended_action") or "").strip().upper()
                
                if "dimensions" in analysis and isinstance(analysis["dimensions"], dict):
                    dimensions = self._validate_dimensions(analysis["dimensions"])
                    self._log(f"Validated {len(dimensions)} dimensions")
                else:
                    # If LLM didn't return dimensions, create defaults
                    self._log("No dimensions in response, using defaults")
                    dimensions = self._validate_dimensions({})

                # Timeliness adjustment using upload dates (S3 + metadata) and content dates
                content_dates_raw = analysis.get("metadata", {}).get("dates", []) if isinstance(analysis.get("metadata", {}), dict) else []
                content_dates = []
                if isinstance(content_dates_raw, list):
                    # Normalize to strings and keep unique order
                    seen = set()
                    for d in content_dates_raw:
                        s = str(d).strip()
                        if s and s not in seen:
                            seen.add(s)
                            content_dates.append(s)

                # Timeliness based only on S3 upload date (as requested)
                if upload_age_days is not None:
                    best_age_days = upload_age_days
                    best_date_iso = upload_date_iso

                    # Always update the canonical Timeliness dimension (capital T) for consistency
                    dim_key = "Timeliness"
                    if "Timeliness" not in dimensions and "timeliness" in dimensions:
                        dim_key = "timeliness"
                    timeliness = dimensions.get(dim_key, {"score": 50, "evidence": "Dimension not assessed by LLM"})

                    if best_age_days > 30:
                        adjusted_score = min(timeliness.get("score", 50), 60)
                        note_parts = [
                            f"Upload date (S3) {best_date_iso} is {best_age_days} days old (>30 days); timeliness reduced to {adjusted_score}.",
                        ]
                        if content_dates:
                            note_parts.append(f"Content dates found: {', '.join(content_dates[:3])}{'...' if len(content_dates)>3 else ''}")
                        existing_evidence = timeliness.get("evidence", "") or "Dimension not assessed by LLM"
                        timeliness["score"] = adjusted_score
                        timeliness["evidence"] = " ".join(note_parts + [existing_evidence]).strip()
                    else:
                        note_parts = [f"Upload date (S3) {best_date_iso} is {best_age_days} days old (<=30 days); timeliness satisfied."]
                        if content_dates:
                            note_parts.append(f"Content dates found: {', '.join(content_dates[:3])}{'...' if len(content_dates)>3 else ''}")
                        else:
                            note_parts.append("No explicit content dates found in document.")
                        existing_evidence = timeliness.get("evidence", "") or "Dimension not assessed by LLM"
                        timeliness["evidence"] = " ".join(note_parts + [existing_evidence]).strip()

                    dimensions[dim_key] = timeliness
                
                # 3.7 Calculate overall quality score from dimensions (average of all 17 dimensions)
                dimension_values = [dim.get("score", 50) for dim in dimensions.values()]
                overall_quality_score = round(sum(dimension_values) / len(dimension_values)) if dimension_values else 50
                self._log(f"Overall quality score calculated: {overall_quality_score}")

                # 3.8 Derive recommended action if LLM did not provide one
                if recommended_action not in {"KEEP", "REVIEW", "QUARANTINE", "DISCARD"}:
                    if overall_quality_score >= 85:
                        recommended_action = "KEEP"
                    elif overall_quality_score >= 70:
                        recommended_action = "REVIEW"
                    elif overall_quality_score >= 60:
                        recommended_action = "QUARANTINE"
                    else:
                        recommended_action = "DISCARD"
                self._log(f"Recommended action: {recommended_action}")

                # 4. Store analysis as individual JSON file
                file_analysis = {
                    "file_key": key,
                    "file_name": file_name,
                    "status": "success",
                    "processed_at": datetime.datetime.utcnow().isoformat() + "Z",
                    "upload_date": upload_date_iso,
                    "upload_age_days": upload_age_days,
                    "embedding": embedding,
                    "overall_quality_score": overall_quality_score,
                    "recommended_action": recommended_action,
                    "bucket": bucket,
                    **analysis,
                    # Ensure our adjusted dimensions override any LLM-provided ones
                    "dimensions": dimensions,
                }
                
                # Write individual JSON file
                json_key = f"{key}.json"
                self._log(f"Writing result to S3: {json_key}")
                self.s3_service.write_file(
                    bucket=bucket,
                    key=json_key,
                    content=json.dumps(file_analysis, indent=2),
                    region=region,
                    access_key=access_key,
                    secret_key=secret_key,
                    role_arn=role_arn
                )
                
                file_analyses.append(file_analysis)
                
                results.append({
                    "file_key": key,
                    "status": "success",
                    "summary": analysis.get("summary", "No summary available"),
                    "upload_date": upload_date_iso,
                    "upload_age_days": upload_age_days,
                    "metadata_key": json_key
                })
                
            except Exception as e:
                self._log(f"Error processing {key}: {str(e)}")
                import traceback
                traceback.print_exc()
                results.append({
                    "file_key": key,
                    "status": "error",
                    "error": str(e)
                })
                file_analyses.append({
                    "file_key": key,
                    "file_name": key.split('/')[-1],
                    "status": "error",
                    "error": str(e),
                    "processed_at": datetime.datetime.utcnow().isoformat() + "Z"
                })

        self._log(f"Processing complete. Processed {len(file_analyses)} files.")
        
        # 4. Calculate cosine similarity for duplicate detection (>95%)
        self._log("=" * 80)
        self._log("STARTING COSINE SIMILARITY CALCULATIONS FOR DUPLICATE DETECTION")
        self._log("=" * 80)
        successful_files = [f for f in file_analyses if f.get("status") == "success"]
        self._log(f"Files eligible for duplicate check: {len(successful_files)}")

        if len(successful_files) < 2:
            self._log("Not enough files to calculate similarity. Need at least 2 files.")
        else:
            total_comparisons = 0
            duplicates_found = 0
            similarity_pairs = []  # track all meta-gated pairs with their cosine for UI

            # Cache summary embeddings to avoid repeated Titan calls
            for f in successful_files:
                f["summary_embedding"] = []

            for i, file1 in enumerate(successful_files):
                potential_duplicates = []
                meta1 = file1.get("metadata", {})
                summary1 = file1.get("summary", "")

                for j, file2 in enumerate(successful_files):
                    if i >= j:  # Skip self and already compared pairs
                        continue

                    meta2 = file2.get("metadata", {})
                    meta_sim = self._metadata_similarity(
                        file1.get("document_type"),
                        file2.get("document_type"),
                        (meta1 or {}).get("topics"),
                        (meta2 or {}).get("topics"),
                        (meta1 or {}).get("key_terms"),
                        (meta2 or {}).get("key_terms"),
                    )
                    total_comparisons += 1

                    self._log(f"Metadata similarity {file1.get('file_name')} <-> {file2.get('file_name')}: {round(meta_sim*100,2)}%")

                    meta_gate_ok = meta_sim >= 0.7

                    # Prefer full-document embedding for duplicate check; fall back to summary if missing.
                    emb1 = file1.get("embedding", [])
                    emb2 = file2.get("embedding", [])

                    # If comprehensive embeddings are missing (e.g., Titan empty), try summary embeddings as fallback
                    if (not emb1 or not emb2):
                        try:
                            if not file1.get("summary_embedding") and summary1:
                                file1["summary_embedding"] = self.bedrock_service.get_embedding(
                                    text=summary1,
                                    region=region,
                                    access_key=access_key,
                                    secret_key=secret_key,
                                    role_arn=role_arn
                                ) or []
                            if not file2.get("summary_embedding") and file2.get("summary"):
                                file2["summary_embedding"] = self.bedrock_service.get_embedding(
                                    text=file2.get("summary"),
                                    region=region,
                                    access_key=access_key,
                                    secret_key=secret_key,
                                    role_arn=role_arn
                                ) or []
                        except Exception as embed_err:
                            self._log(f"Summary embedding failed: {str(embed_err)}")

                        emb1 = emb1 or file1.get("summary_embedding", [])
                        emb2 = emb2 or file2.get("summary_embedding", [])

                    # Fallback: bag-of-words embedding if Bedrock embeddings are empty
                    if (not emb1 or not emb2):
                        import re
                        from collections import Counter

                        def bow_counts(text: str) -> Counter:
                            tokens = [t for t in re.split(r"\W+", text.lower()) if t]
                            return Counter(tokens)

                        c1 = bow_counts(summary1 or file1.get("context", ""))
                        c2 = bow_counts(file2.get("summary", "") or file2.get("context", ""))

                        vocab = sorted(set(c1.keys()) | set(c2.keys()))
                        if vocab:
                            # Align vocab across both vectors to ensure cosine works
                            emb1 = emb1 or [c1.get(k, 0) for k in vocab]
                            emb2 = emb2 or [c2.get(k, 0) for k in vocab]

                    if not emb1 or not emb2:
                        self._log(f"Skipping pair due to missing embeddings: {file1.get('file_name')} / {file2.get('file_name')}")
                        continue

                    similarity = self._cosine_similarity(emb1, emb2)
                    self._log(f"Summary cosine {file1.get('file_name')} <-> {file2.get('file_name')}: {round(similarity*100,2)}% (meta {round(meta_sim*100,2)}%)")

                    similarity_pairs.append({
                        "file_1": file1.get("file_name"),
                        "file_2": file2.get("file_name"),
                        "similarity": round(similarity * 100, 2),
                        "metadata_similarity": round(meta_sim * 100, 2)
                    })

                    if meta_gate_ok and similarity >= 0.95:
                        duplicates_found += 1
                        entry = {
                            "file_name": file2.get("file_name"),
                            "file_key": file2.get("file_key"),
                            "similarity": round(similarity * 100, 2),
                            "metadata_similarity": round(meta_sim * 100, 2)
                        }
                        potential_duplicates.append(entry)
                        # Symmetric add to file2 as well
                        other_list = file2.get("potential_duplicates", [])
                        other_list.append({
                            "file_name": file1.get("file_name"),
                            "file_key": file1.get("file_key"),
                            "similarity": round(similarity * 100, 2),
                            "metadata_similarity": round(meta_sim * 100, 2)
                        })
                        file2["potential_duplicates"] = other_list
                        self._log(f"⚠️  DUPLICATE DETECTED: {file1.get('file_name')} <-> {file2.get('file_name')} ({round(similarity * 100, 2)}%)")

                if potential_duplicates:
                    potential_duplicates.sort(key=lambda x: x["similarity"], reverse=True)
                    file1["potential_duplicates"] = potential_duplicates
                    self._log(f"File '{file1.get('file_name')}' has {len(potential_duplicates)} potential duplicate(s)")

            self._log(f"Duplicate detection complete. Total comparisons: {total_comparisons}, Duplicates found: {duplicates_found}")
            # Attach all similarity pairs for UI transparency
            similarity_pairs.sort(key=lambda x: x.get("similarity", 0), reverse=True)
            for f in file_analyses:
                f["similarity_pairs"] = similarity_pairs
            # Also attach to lightweight results list so UI can read without consolidated_json
            for r in results:
                r["similarity_pairs"] = similarity_pairs
        
        self._log("=" * 80)
        
        # Return results with full analysis data for immediate UI display
        # We attach the full analysis list to the first result so the UI can grab it easily
        # This maintains backward compatibility with the UI we just built
        
        consolidated_json = {
            "processed_at": datetime.datetime.utcnow().isoformat() + "Z",
            "total_files": len(file_keys),
            "successful": len([r for r in results if r["status"] == "success"]),
            "failed": len([r for r in results if r["status"] == "error"]),
            "model_used": model_id,
            "files": file_analyses
        }

        # Add a flat list of duplicate pairs for easy UI display
        duplicate_pairs = []
        seen_pairs = set()
        for f in file_analyses:
            for dup in f.get("potential_duplicates", []) or []:
                a = f.get("file_name")
                b = dup.get("file_name")
                if not a or not b:
                    continue
                key = tuple(sorted([a, b]))
                if key in seen_pairs:
                    continue
                seen_pairs.add(key)
                duplicate_pairs.append({
                    "file_1": a,
                    "file_2": b,
                    "similarity": dup.get("similarity"),
                    "metadata_similarity": dup.get("metadata_similarity")
                })

        duplicate_pairs.sort(key=lambda x: x.get("similarity", 0), reverse=True)
        consolidated_json["duplicate_pairs"] = duplicate_pairs

        # Add full similarity pairs list for UI when many files are selected
        if 'similarity_pairs' in locals():
            consolidated_json["similarity_pairs"] = similarity_pairs

        # Mirror summary pairs onto first result for backward compatibility with UI
        if results:
            results[0]["similarity_pairs"] = similarity_pairs
        
        # Save consolidated JSON to S3 in output_folder
        timestamp = datetime.datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        consolidated_key = f"output_folder/quality_check_results_{timestamp}.json"
        self._log(f"Saving consolidated results to S3: {consolidated_key}")
        
        try:
            await self.s3_service.write_file(
                bucket=bucket,
                key=consolidated_key,
                content=json.dumps(consolidated_json, indent=2),
                region=region,
                access_key=access_key,
                secret_key=secret_key,
                role_arn=role_arn
            )
            self._log(f"Consolidated results saved successfully to {consolidated_key}")
        except Exception as e:
            self._log(f"Warning: Failed to save consolidated results: {str(e)}")
            
        # Save LOCALLY as requested by user
        try:
            import os
            local_dir = "data/results"
            os.makedirs(local_dir, exist_ok=True)
            local_filename = f"{local_dir}/results_{bucket}_{timestamp}.json"
            with open(local_filename, "w") as f:
                json.dump(consolidated_json, f, indent=2)
            self._log(f"Saved local result copy to {local_filename}")
        except Exception as e:
            self._log(f"Failed to save local result: {str(e)}")
        
        if results:
            results[0]["consolidated_json"] = consolidated_json
            results[0]["consolidated_key"] = consolidated_key
            
        return results

    def list_local_history(self) -> List[Dict[str, Any]]:
        """List all locally saved result files"""
        import os
        import glob
        
        local_dir = "data/results"
        if not os.path.exists(local_dir):
            return []
            
        files = []
        for filepath in glob.glob(f"{local_dir}/*.json"):
            try:
                filename = os.path.basename(filepath)
                stats = os.stat(filepath)
                with open(filepath, 'r') as f:
                    data = json.load(f)
                    
                files.append({
                    "filename": filename,
                    "created_at": datetime.datetime.fromtimestamp(stats.st_mtime).isoformat(),
                    "total_files": data.get("total_files", 0),
                    "successful": data.get("successful", 0),
                    "failed": data.get("failed", 0),
                    "model_used": data.get("model_used", "unknown")
                })
            except Exception as e:
                self._log(f"Error reading local file {filepath}: {str(e)}")
                
        # Sort by creation time (newest first)
        files.sort(key=lambda x: x['created_at'], reverse=True)
        return files

    def get_local_history_content(self, filename: str) -> Dict[str, Any]:
        """Get content of a specific local result file"""
        import os
        local_dir = "data/results"
        filepath = os.path.join(local_dir, filename)
        
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"Local file not found: {filename}")
            
        with open(filepath, 'r') as f:
            return json.load(f)

    def get_all_local_results(self) -> List[Dict[str, Any]]:
        """Get all local results with full content (including embeddings)"""
        import os
        import glob
        
        local_dir = "data/results"
        if not os.path.exists(local_dir):
            return []
            
        results = []
        for filepath in glob.glob(f"{local_dir}/*.json"):
            try:
                with open(filepath, 'r') as f:
                    data = json.load(f)
                    # Flatten if it's a consolidated result
                    if 'files' in data and isinstance(data['files'], list):
                        results.extend(data['files'])
                    else:
                        results.append(data)
            except Exception:
                continue
        return results

    async def get_file_content(self, bucket: str, key: str, region: str = None, access_key: str = None, secret_key: str = None, role_arn: str = None) -> Dict[str, Any]:
        self._log(f"get_file_content called for bucket={bucket}, key={key}")
        try:
            content_bytes = await self.s3_service.read_file(bucket, key, region, access_key, secret_key, role_arn)
            return json.loads(content_bytes)
        except FileNotFoundError:
            # Auto-reconstruction: If the summary file is missing, try to build it from individual files
            if "quality_check_results" in key:
                self._log(f"Summary file {key} not found. Attempting to reconstruct from individual files...")
                try:
                    # List files in root AND output_folder to be sure we catch everything
                    all_files = []
                    
                    # 1. Check root
                    root_data = await self.s3_service.list_files(bucket, "", region, access_key, secret_key, role_arn)
                    all_files.extend(root_data.get('files', []))
                    
                    # 2. Check output_folder
                    out_data = await self.s3_service.list_files(bucket, "output_folder/", region, access_key, secret_key, role_arn)
                    all_files.extend(out_data.get('files', []))
                    
                    # Filter for individual analysis files
                    json_files = [f for f in all_files if f.get('key', '').endswith('.json') and 
                                  "quality_check_results" not in f.get('key', '') and 
                                  not f.get('is_folder', False)]
                    
                    # Remove duplicates based on key
                    unique_files = {f['key']: f for f in json_files}.values()
                    json_files = list(unique_files)
                    
                    # Sort by most recent
                    json_files.sort(key=lambda x: x.get('last_modified', ''), reverse=True)
                    json_files = json_files[:50]
                    
                    reconstructed_files = []
                    successful = 0
                    failed = 0
                    
                    # Add a DEBUG entry so the user sees SOMETHING
                    reconstructed_files.append({
                        "file_name": "🔍 SYSTEM DIAGNOSTIC",
                        "status": "info",
                        "summary": f"Found {len(json_files)} potential analysis files. Attempting to load...",
                        "processed_at": datetime.datetime.utcnow().isoformat() + "Z"
                    })
                    
                    for file_info in json_files:
                        try:
                            content = await self.s3_service.read_file(bucket, file_info['key'], region, access_key, secret_key, role_arn)
                            data = json.loads(content)
                            reconstructed_files.append(data)
                            if data.get('status') == 'success':
                                successful += 1
                            else:
                                failed += 1
                        except Exception as read_err:
                            reconstructed_files.append({
                                "file_name": file_info['key'],
                                "status": "error",
                                "error": f"Failed to read: {str(read_err)}",
                                "processed_at": datetime.datetime.utcnow().isoformat() + "Z"
                            })
                            
                    # Construct consolidated response
                    consolidated_data = {
                        "processed_at": datetime.datetime.utcnow().isoformat() + "Z",
                        "total_files": len(reconstructed_files),
                        "successful": successful,
                        "failed": failed,
                        "model_used": "reconstructed",
                        "files": reconstructed_files
                    }
                    
                    self._log("Successfully reconstructed results")
                    return consolidated_data
                    
                except Exception as reconstruct_err:
                    self._log(f"Reconstruction failed: {str(reconstruct_err)}")
                    raise FileNotFoundError(f"File not found: {key} and failed to reconstruct.")
            
            raise
        except Exception as e:
            self._log(f"Error reading file content {key}: {str(e)}")
            raise e
    
    async def reconstruct_results(self, bucket: str, region: str = None, access_key: str = None, secret_key: str = None, role_arn: str = None) -> Dict[str, Any]:
        """
        Explicitly reconstruct results by scanning the bucket for all analysis files.
        This is used by the 'Load Past Results' button.
        """
        self._log(f"reconstruct_results called for bucket={bucket}")
        try:
            # List files in root AND output_folder
            all_files = []
            
            # 1. Check root
            root_data = await self.s3_service.list_files(bucket, "", region, access_key, secret_key, role_arn)
            all_files.extend(root_data.get('files', []))
            
            # 2. Check output_folder
            out_data = await self.s3_service.list_files(bucket, "output_folder/", region, access_key, secret_key, role_arn)
            all_files.extend(out_data.get('files', []))
            
            # Filter for individual analysis files
            json_files = [f for f in all_files if f.get('key', '').endswith('.json') and 
                          "quality_check_results" not in f.get('key', '') and 
                          not f.get('is_folder', False)]
            
            # Remove duplicates
            unique_files = {f['key']: f for f in json_files}.values()
            json_files = list(unique_files)
            
            # Sort by most recent
            json_files.sort(key=lambda x: x.get('last_modified', ''), reverse=True)
            json_files = json_files[:50]
            
            reconstructed_files = []
            successful = 0
            failed = 0
            
            for file_info in json_files:
                try:
                    content = await self.s3_service.read_file(bucket, file_info['key'], region, access_key, secret_key, role_arn)
                    data = json.loads(content)
                    reconstructed_files.append(data)
                    if data.get('status') == 'success':
                        successful += 1
                    else:
                        failed += 1
                except Exception as read_err:
                    reconstructed_files.append({
                        "file_name": file_info['key'],
                        "status": "error",
                        "error": f"Failed to read: {str(read_err)}",
                        "processed_at": datetime.datetime.utcnow().isoformat() + "Z"
                    })
            
            if not reconstructed_files:
                # Return empty structure instead of error
                return {
                    "processed_at": datetime.datetime.utcnow().isoformat() + "Z",
                    "total_files": 0,
                    "successful": 0,
                    "failed": 0,
                    "model_used": "none",
                    "files": [],
                    "note": "No analysis files found in bucket."
                }
                
            return {
                "processed_at": datetime.datetime.utcnow().isoformat() + "Z",
                "total_files": len(reconstructed_files),
                "successful": successful,
                "failed": failed,
                "model_used": "reconstructed",
                "files": reconstructed_files,
                "note": "Loaded from history"
            }
            
        except Exception as e:
            self._log(f"Error reconstructing results: {str(e)}")
            raise e

    async def get_scan_history(self, bucket: str, prefix: str = "", region: str = None, access_key: str = None, secret_key: str = None, role_arn: str = None, limit: int = 10) -> List[Dict[str, Any]]:
        """
        Fetch recent scan history by reading .json result files from S3
        Returns aggregated scan data for dashboard visualization
        """
        self._log(f"get_scan_history called for bucket={bucket}, prefix={prefix}")
        try:
            # List all files in the bucket
            files_data = await self.s3_service.list_files(bucket, prefix, region, access_key, secret_key, role_arn)
            files = files_data.get('files', [])
            
            # Filter for .json files (our scan results)
            json_files = [f for f in files if f.get('key', '').endswith('.json') and not f.get('is_folder', False)]
            
            # Sort by last_modified (most recent first)
            json_files.sort(key=lambda x: x.get('last_modified', ''), reverse=True)
            
            # Limit to requested number
            json_files = json_files[:limit]
            
            # Fetch content of each JSON file
            scan_results = []
            for file_info in json_files:
                try:
                    content = await self.get_file_content(bucket, file_info['key'], region, access_key, secret_key, role_arn)
                    
                    # Extract relevant information
                    scan_entry = {
                        'file_key': file_info['key'],
                        'file_name': content.get('file_name', file_info['key']),
                        'processed_at': content.get('processed_at', file_info.get('last_modified', '')),
                        'quality_score': content.get('quality_score', 0),
                        'status': content.get('status', 'unknown'),
                        'summary': content.get('summary', ''),
                        'source_file': content.get('file_key', '').replace('.json', ''),
                    }
                    scan_results.append(scan_entry)
                except Exception as e:
                    self._log(f"Error reading scan result {file_info['key']}: {str(e)}")
                    continue
            
            # Calculate aggregate statistics
            total_scans = len(scan_results)
            successful_scans = len([s for s in scan_results if s['status'] == 'success'])
            avg_quality = sum(s['quality_score'] for s in scan_results) / total_scans if total_scans > 0 else 0
            
            return {
                'scans': scan_results,
                'total_scans': total_scans,
                'successful_scans': successful_scans,
                'failed_scans': total_scans - successful_scans,
                'average_quality_score': round(avg_quality, 1)
            }
            
        except Exception as e:
            self._log(f"Error fetching scan history: {str(e)}")
            raise e
